#!/usr/bin/env python3
"""Run deterministic regression tests for the bundled presentation skills."""

from __future__ import annotations

import argparse
import base64
import io
import shutil
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock

from defusedxml import minidom
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]


def companion_root(name: str) -> Path:
    bundled = REPO_ROOT / "companion" / name
    installed = REPO_ROOT.parent / name
    if bundled.is_dir():
        return bundled
    if installed.is_dir():
        return installed
    raise RuntimeError(f"required companion skill is unavailable: {name}")


PPTX_OPERATOR = companion_root("pptx-operator")
SVG_ENGINE = companion_root("svg-diagram-engine")
PPTX_SCRIPTS = PPTX_OPERATOR / "scripts"
SVG_SCRIPTS = SVG_ENGINE / "scripts"
SVG_EXAMPLES = SVG_ENGINE / "assets/examples"
IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_EXTENSION_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"


ARTIFACT_ROOT: Path
RENDER_BACKEND = "none"


def configure_utf8_stdio() -> None:
    """Make Windows test output deterministic for Chinese paths and slide text."""
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if reconfigure is not None:
            try:
                reconfigure(encoding="utf-8", errors="backslashreplace")
            except (AttributeError, ValueError):
                pass


configure_utf8_stdio()


def run_python(script: Path, *arguments: object) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(script), *(str(argument) for argument in arguments)],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=False,
    )


def command_output(result: subprocess.CompletedProcess[str]) -> str:
    return "\n".join(part for part in (result.stdout.strip(), result.stderr.strip()) if part)


def direct_child(node, tag_name: str):
    for child in node.childNodes:
        if child.nodeType == child.ELEMENT_NODE and child.tagName == tag_name:
            return child
    return None


def solid_png(path: Path, color: tuple[int, int, int]) -> bytes:
    image = Image.new("RGB", (96, 64), "white")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((4, 4, 92, 60), radius=12, fill=color)
    image.save(path, format="PNG")
    return path.read_bytes()


def svg_with_image(png_bytes: bytes, external_href: str | None = None) -> bytes:
    href = external_href or (
        "data:image/png;base64," + base64.b64encode(png_bytes).decode("ascii")
    )
    return f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 320 180">
  <title>Embedded image portability fixture</title>
  <rect width="320" height="180" rx="20" fill="#F4F8FF"/>
  <path d="M32 90 H140" stroke="#1677FF" stroke-width="6"/>
  <image x="180" y="58" width="96" height="64" href="{href}"/>
</svg>'''.encode("utf-8")


def create_base_media_deck(png_path: Path, output: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), Inches(2), Inches(1.5), width=Inches(2.4))
    presentation.save(output)


def add_svg_alternative(base_deck: Path, svg_bytes: bytes, output: Path) -> None:
    slide_part = "ppt/slides/slide1.xml"
    rel_part = "ppt/slides/_rels/slide1.xml.rels"
    content_types_part = "[Content_Types].xml"
    svg_media_part = "ppt/media/harness-diagram.svg"

    with zipfile.ZipFile(base_deck) as source:
        slide_document = minidom.parseString(source.read(slide_part))
        relationship_document = minidom.parseString(source.read(rel_part))
        content_types_document = minidom.parseString(source.read(content_types_part))

        relationship_ids = {
            element.getAttribute("Id")
            for element in relationship_document.getElementsByTagName("Relationship")
        }
        index = 1
        while f"rId{index}" in relationship_ids:
            index += 1
        svg_relationship_id = f"rId{index}"

        relationship = relationship_document.createElement("Relationship")
        relationship.setAttribute("Id", svg_relationship_id)
        relationship.setAttribute("Type", IMAGE_REL)
        relationship.setAttribute("Target", "../media/harness-diagram.svg")
        relationship_document.documentElement.appendChild(relationship)

        blips = slide_document.getElementsByTagName("a:blip")
        if not blips:
            raise RuntimeError("media fixture has no DrawingML image blip")
        extension_list = slide_document.createElement("a:extLst")
        extension = slide_document.createElement("a:ext")
        extension.setAttribute("uri", SVG_EXTENSION_URI)
        svg_blip = slide_document.createElement("asvg:svgBlip")
        svg_blip.setAttribute("xmlns:asvg", SVG_NS)
        svg_blip.setAttribute("r:embed", svg_relationship_id)
        extension.appendChild(svg_blip)
        extension_list.appendChild(extension)
        blips[0].appendChild(extension_list)

        defaults = content_types_document.getElementsByTagName("Default")
        if not any(item.getAttribute("Extension").casefold() == "svg" for item in defaults):
            default = content_types_document.createElement("Default")
            default.setAttribute("Extension", "svg")
            default.setAttribute("ContentType", "image/svg+xml")
            content_types_document.documentElement.appendChild(default)

        with zipfile.ZipFile(output, "w") as destination:
            for info in source.infolist():
                data = source.read(info.filename)
                if info.filename == slide_part:
                    data = slide_document.toxml(encoding="UTF-8")
                elif info.filename == rel_part:
                    data = relationship_document.toxml(encoding="UTF-8")
                elif info.filename == content_types_part:
                    data = content_types_document.toxml(encoding="UTF-8")
                destination.writestr(info, data)
            destination.writestr(svg_media_part, svg_bytes)


def replace_zip_part(source_path: Path, output_path: Path, part: str, data: bytes) -> None:
    with zipfile.ZipFile(source_path) as source, zipfile.ZipFile(output_path, "w") as destination:
        for info in source.infolist():
            destination.writestr(info, data if info.filename == part else source.read(info.filename))


def make_external_relationship(source_path: Path, output_path: Path) -> None:
    part = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(source_path) as source:
        document = minidom.parseString(source.read(part))
    for relationship in document.getElementsByTagName("Relationship"):
        if relationship.getAttribute("Target").endswith("harness-diagram.svg"):
            relationship.setAttribute("Target", "file:///tmp/harness-diagram.svg")
            relationship.setAttribute("TargetMode", "External")
            break
    else:
        raise RuntimeError("SVG relationship was not found")
    replace_zip_part(source_path, output_path, part, document.toxml(encoding="UTF-8"))


def make_negative_connector_deck(output: Path) -> None:
    normalized = output.with_name("connector-normalized-source.pptx")
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1),
        Inches(4),
        Inches(5),
        Inches(1),
    )
    presentation.save(normalized)

    part = "ppt/slides/slide1.xml"
    with zipfile.ZipFile(normalized) as source:
        document = minidom.parseString(source.read(part))
    connectors = document.getElementsByTagName("p:cxnSp")
    if not connectors:
        raise RuntimeError("connector fixture has no connector")
    shape_properties = direct_child(connectors[0], "p:spPr")
    transform = direct_child(shape_properties, "a:xfrm")
    offset = direct_child(transform, "a:off")
    extent = direct_child(transform, "a:ext")
    height = int(extent.getAttribute("cy"))
    offset.setAttribute("y", str(int(offset.getAttribute("y")) + height))
    extent.setAttribute("cy", str(-height))
    transform.removeAttribute("flipV")
    replace_zip_part(normalized, output, part, document.toxml(encoding="UTF-8"))


def create_arial_chinese_deck(output: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "第五课 Transformer 推理性能"
    run.font.name = "Arial"
    run.font.size = Pt(30)
    presentation.save(output)


class SkillHarness(unittest.TestCase):
    maxDiff = None

    @classmethod
    def setUpClass(cls) -> None:
        cls.root = ARTIFACT_ROOT
        cls.assets = cls.root / "source-assets"
        cls.assets.mkdir(parents=True)
        fallback_path = cls.assets / "fallback.png"
        png_bytes = solid_png(fallback_path, (22, 119, 255))
        cls.good_svg_bytes = svg_with_image(png_bytes)
        cls.external_svg_bytes = svg_with_image(png_bytes, "file:///tmp/source-image.png")
        cls.invalid_svg_bytes = svg_with_image(
            png_bytes, "data:image/png;base64,SGVsbG8="
        )

        base_deck = cls.root / "media-base.pptx"
        cls.good_deck = cls.root / "embedded-media.pptx"
        create_base_media_deck(fallback_path, base_deck)
        add_svg_alternative(base_deck, cls.good_svg_bytes, cls.good_deck)
        cls.external_relationship_deck = cls.root / "external-relationship.pptx"
        make_external_relationship(cls.good_deck, cls.external_relationship_deck)
        cls.external_svg_deck = cls.root / "external-svg-reference.pptx"
        replace_zip_part(
            cls.good_deck,
            cls.external_svg_deck,
            "ppt/media/harness-diagram.svg",
            cls.external_svg_bytes,
        )
        cls.invalid_svg_deck = cls.root / "invalid-svg-image-bytes.pptx"
        replace_zip_part(
            cls.good_deck,
            cls.invalid_svg_deck,
            "ppt/media/harness-diagram.svg",
            cls.invalid_svg_bytes,
        )
        shutil.rmtree(cls.assets)

    def assert_success(self, result: subprocess.CompletedProcess[str]) -> str:
        output = command_output(result)
        self.assertEqual(result.returncode, 0, output)
        return output

    def assert_failure(self, result: subprocess.CompletedProcess[str]) -> str:
        output = command_output(result)
        self.assertNotEqual(result.returncode, 0, output)
        return output

    def test_01_svg_embedded_image_and_renderer(self) -> None:
        good_svg = self.root / "embedded-image.svg"
        bad_svg = self.root / "external-image.svg"
        invalid_svg = self.root / "invalid-image-bytes.svg"
        rendered_png = self.root / "embedded-image.png"
        good_svg.write_bytes(self.good_svg_bytes)
        bad_svg.write_bytes(self.external_svg_bytes)
        invalid_svg.write_bytes(self.invalid_svg_bytes)

        self.assert_success(run_python(SVG_SCRIPTS / "svg_validate.py", good_svg, "--strict"))
        bad_output = self.assert_failure(
            run_python(SVG_SCRIPTS / "svg_validate.py", bad_svg, "--strict")
        )
        self.assertIn("external href is not portable", bad_output)
        invalid_output = self.assert_failure(
            run_python(SVG_SCRIPTS / "svg_validate.py", invalid_svg, "--strict")
        )
        self.assertIn("do not match the declared media type", invalid_output)

        render_result = run_python(SVG_SCRIPTS / "svg_render.py", good_svg, rendered_png)
        if render_result.returncode and "no non-browser SVG backend" in command_output(render_result):
            self.skipTest("no SVG renderer is installed")
        self.assert_success(render_result)
        with Image.open(rendered_png) as image:
            self.assertGreater(image.width, 0)
            self.assertGreater(image.height, 0)

    def test_02_structured_svg_smoke(self) -> None:
        svg_path = self.root / "structured-task-tree.svg"
        self.assert_success(
            run_python(
                SVG_SCRIPTS / "diagram_render.py",
                SVG_EXAMPLES / "task-tree.json",
                svg_path,
            )
        )
        self.assert_success(run_python(SVG_SCRIPTS / "svg_validate.py", svg_path))

    def test_03_pptx_svg_and_png_are_internal(self) -> None:
        moved_directory = self.root / "moved/deep/location"
        moved_directory.mkdir(parents=True)
        moved_deck = moved_directory / "portable.pptx"
        shutil.copy2(self.good_deck, moved_deck)

        validation = self.assert_success(
            run_python(PPTX_SCRIPTS / "validate_pptx.py", moved_deck)
        )
        audit = self.assert_success(run_python(PPTX_SCRIPTS / "audit_media.py", moved_deck))
        self.assertIn("embedded_svg_blips=1", audit)
        self.assertIn("svg_files=1", audit)
        self.assertIn("svg_embedded_images=1", audit)
        self.assertIn("svg_external_references=0", audit)
        self.assertIn("PASS", validation)

        with zipfile.ZipFile(moved_deck) as package:
            media = {name for name in package.namelist() if name.startswith("ppt/media/")}
            self.assertTrue(any(name.endswith(".png") for name in media), media)
            self.assertIn("ppt/media/harness-diagram.svg", media)
            embedded_svg = package.read("ppt/media/harness-diagram.svg")
            self.assertIn(b"data:image/png;base64,", embedded_svg)
            self.assertNotIn(b"file://", embedded_svg)

        if RENDER_BACKEND != "none":
            render_dir = self.root / "moved-render"
            self.assert_success(
                run_python(
                    PPTX_SCRIPTS / "render_pptx.py",
                    moved_deck,
                    "--output-dir",
                    render_dir,
                    "--backend",
                    RENDER_BACKEND,
                )
            )

    def test_04_external_media_is_rejected(self) -> None:
        relationship_output = self.assert_failure(
            run_python(PPTX_SCRIPTS / "audit_media.py", self.external_relationship_deck)
        )
        self.assertIn("external non-hyperlink relationship", relationship_output)

        svg_output = self.assert_failure(
            run_python(PPTX_SCRIPTS / "audit_media.py", self.external_svg_deck)
        )
        self.assertIn("external reference inside ppt/media/harness-diagram.svg", svg_output)
        invalid_svg_output = self.assert_failure(
            run_python(PPTX_SCRIPTS / "audit_media.py", self.invalid_svg_deck)
        )
        self.assertIn("invalid embedded image bytes", invalid_svg_output)

    def test_05_negative_connector_is_rejected_and_repaired(self) -> None:
        invalid = self.root / "negative-connector.pptx"
        fixed = self.root / "normalized-connector.pptx"
        make_negative_connector_deck(invalid)
        invalid_output = self.assert_failure(
            run_python(PPTX_SCRIPTS / "validate_pptx.py", invalid)
        )
        self.assertIn("negative transform extent", invalid_output)
        self.assertIn("kind=connector", invalid_output)

        self.assert_success(
            run_python(PPTX_SCRIPTS / "normalize_connectors.py", invalid, fixed)
        )
        self.assert_success(run_python(PPTX_SCRIPTS / "validate_pptx.py", fixed))
        with zipfile.ZipFile(fixed) as package:
            xml = package.read("ppt/slides/slide1.xml")
            self.assertIn(b'flipV="1"', xml)
            self.assertNotIn(b'cy="-', xml)

    def test_06_chinese_font_runs_are_normalized(self) -> None:
        source = self.root / "arial-chinese.pptx"
        fixed = self.root / "explicit-cjk.pptx"
        create_arial_chinese_deck(source)
        failure = self.assert_failure(
            run_python(
                PPTX_SCRIPTS / "audit_pptx_fonts.py",
                source,
                "--libreoffice-safe",
                "--strict",
            )
        )
        self.assertIn("Arial", failure)
        self.assert_success(
            run_python(
                PPTX_SCRIPTS / "apply_cjk_fonts.py",
                source,
                fixed,
                "--title-font",
                "Harness CJK",
                "--body-font",
                "Harness CJK",
                "--allow-uninstalled",
            )
        )
        self.assert_success(
            run_python(
                PPTX_SCRIPTS / "audit_pptx_fonts.py",
                fixed,
                "--libreoffice-safe",
                "--strict",
            )
        )

    def test_07_windows_font_aliases(self) -> None:
        sys.path.insert(0, str(PPTX_SCRIPTS))
        try:
            import font_policy

            self.assertEqual(
                font_policy._pick(["宋体"], {"simsun": "SimSun"})[0],
                "SimSun",
            )
            self.assertEqual(
                font_policy._pick(
                    ["Microsoft YaHei"], {"微软雅黑": "微软雅黑"}
                )[0],
                "微软雅黑",
            )
            self.assertEqual(
                font_policy._pick(["等线"], {"dengxian": "DengXian"})[0],
                "DengXian",
            )
        finally:
            sys.path.pop(0)

    def test_08_source_han_is_default_cjk_family(self) -> None:
        sys.path.insert(0, str(PPTX_SCRIPTS))
        try:
            import font_policy

            installed = {
                "source han sans sc": "Source Han Sans SC",
                "noto sans cjk sc": "Noto Sans CJK SC",
                "pingfang sc": "PingFang SC",
                "microsoft yahei": "Microsoft YaHei",
            }
            for candidates in font_policy.CJK_PREFERENCES.values():
                self.assertEqual(
                    font_policy._pick(candidates, installed)[0],
                    "Source Han Sans SC",
                )
            self.assertEqual(
                font_policy._pick(font_policy.LATIN_PREFERENCES, installed)[0],
                "Source Han Sans SC",
            )
        finally:
            sys.path.pop(0)

    def test_09_font_license_audit_fails_closed(self) -> None:
        source = self.root / "proprietary-font.pptx"
        fixed = self.root / "open-font.pptx"
        create_arial_chinese_deck(source)
        failure = self.assert_failure(
            run_python(PPTX_SCRIPTS / "audit_font_licenses.py", source, "--strict")
        )
        self.assertIn("no verified commercial-use and redistribution license", failure)
        self.assert_success(
            run_python(
                PPTX_SCRIPTS / "apply_cjk_fonts.py",
                source,
                fixed,
                "--title-font",
                "Source Han Sans SC",
                "--body-font",
                "Source Han Sans SC",
                "--allow-uninstalled",
            )
        )
        self.assert_success(
            run_python(PPTX_SCRIPTS / "audit_font_licenses.py", fixed, "--strict")
        )

    def test_10_svg_font_license_audit_fails_closed(self) -> None:
        proprietary = self.root / "proprietary-font.svg"
        approved = self.root / "open-font.svg"
        template = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200">'
            '<title>Font license fixture</title>'
            '<text x="20" y="80" font-family="{font}" font-size="24">PPT 字体审计</text>'
            '</svg>'
        )
        proprietary.write_text(template.format(font="Arial, sans-serif"), encoding="utf-8")
        approved.write_text(
            template.format(font="Source Han Sans SC, sans-serif"), encoding="utf-8"
        )
        failure = self.assert_failure(
            run_python(SVG_SCRIPTS / "svg_validate.py", proprietary, "--strict")
        )
        self.assertIn("no verified commercial-use and redistribution license", failure)
        self.assert_success(run_python(SVG_SCRIPTS / "svg_validate.py", approved, "--strict"))

    def test_11_windows_powerpoint_com_export_contract(self) -> None:
        sys.path.insert(0, str(PPTX_SCRIPTS))
        try:
            import render_pptx

            deck = self.root / "Windows 中文路径.pptx"
            pdf = self.root / "Windows 中文路径.pdf"
            deck.write_bytes(b"pptx-fixture")

            def fake_run(command, **kwargs):
                self.assertIn("-File", command)
                script_path = Path(command[command.index("-File") + 1])
                script_bytes = script_path.read_bytes()
                self.assertTrue(script_bytes.startswith(b"\xef\xbb\xbf"))
                script = script_bytes.decode("utf-8-sig")
                self.assertIn("New-Object -ComObject PowerPoint.Application", script)
                self.assertIn("$presentation.SaveAs($OutputPath, 32)", script)
                self.assertIn("$powerPoint.Quit()", script)
                output = Path(command[command.index("-OutputPath") + 1])
                output.write_bytes(b"%PDF-1.7\n")
                return subprocess.CompletedProcess(command, 0, stdout="")

            with mock.patch.object(render_pptx.subprocess, "run", side_effect=fake_run):
                render_pptx.export_with_powerpoint_windows(
                    deck, pdf, powershell="powershell.exe"
                )
            self.assertTrue(pdf.is_file())

            render_pptx.has_powerpoint.cache_clear()
            with (
                mock.patch.object(render_pptx.sys, "platform", "win32"),
                mock.patch.object(render_pptx, "find_powershell", return_value="powershell.exe"),
                mock.patch.object(
                    render_pptx.subprocess,
                    "run",
                    return_value=subprocess.CompletedProcess([], 0),
                ),
            ):
                self.assertTrue(render_pptx.has_powerpoint())
            render_pptx.has_powerpoint.cache_clear()
        finally:
            sys.path.pop(0)

    def test_12_windows_powerpoint_svg_export_contract(self) -> None:
        sys.path.insert(0, str(SVG_SCRIPTS))
        try:
            import svg_render

            source = self.root / "Windows SVG 中文.svg"
            output = self.root / "Windows SVG 中文.png"
            source.write_text(
                '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200"/>',
                encoding="utf-8",
            )

            def fake_run(command, **kwargs):
                script_path = Path(command[command.index("-File") + 1])
                script_bytes = script_path.read_bytes()
                self.assertTrue(script_bytes.startswith(b"\xef\xbb\xbf"))
                script = script_bytes.decode("utf-8-sig")
                self.assertIn("$slide.Shapes.AddPicture", script)
                self.assertIn('$slide.Export($OutputPath, "PNG"', script)
                destination = Path(command[command.index("-OutputPath") + 1])
                destination.write_bytes(b"\x89PNG\r\n\x1a\n")
                return subprocess.CompletedProcess(command, 0, stdout="")

            with mock.patch.object(svg_render.subprocess, "run", side_effect=fake_run):
                svg_render.render_with_windows_powerpoint(
                    source,
                    output,
                    400,
                    200,
                    1.5,
                    powershell="powershell.exe",
                )
            self.assertTrue(output.is_file())
        finally:
            sys.path.pop(0)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--render-backend",
        choices=("none", "auto", "powerpoint", "libreoffice"),
        default="none",
        help="optionally render the moved PPTX with a real backend",
    )
    parser.add_argument(
        "--artifacts-dir",
        type=Path,
        help="keep generated fixtures and reports under a new subdirectory here",
    )
    args = parser.parse_args()

    global ARTIFACT_ROOT, RENDER_BACKEND
    RENDER_BACKEND = args.render_backend
    temporary = None
    if args.artifacts_dir:
        parent = args.artifacts_dir.expanduser().resolve()
        parent.mkdir(parents=True, exist_ok=True)
        ARTIFACT_ROOT = Path(tempfile.mkdtemp(prefix="pptx-skill-harness-", dir=parent))
    else:
        temporary = tempfile.TemporaryDirectory(prefix="pptx-skill-harness-")
        ARTIFACT_ROOT = Path(temporary.name)

    print(f"artifacts={ARTIFACT_ROOT}")
    print(f"render_backend={RENDER_BACKEND}")
    suite = unittest.defaultTestLoader.loadTestsFromTestCase(SkillHarness)
    result = unittest.TextTestRunner(verbosity=2).run(suite)
    if args.artifacts_dir:
        print(f"artifacts_kept={ARTIFACT_ROOT}")
    if temporary is not None:
        temporary.cleanup()
    raise SystemExit(0 if result.wasSuccessful() else 1)


if __name__ == "__main__":
    main()
