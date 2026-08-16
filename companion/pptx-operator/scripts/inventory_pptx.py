#!/usr/bin/env python3
"""Print a recursive structural inventory for a PPTX or POTX file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from cli_compat import configure_utf8_stdio


configure_utf8_stdio()

def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck", type=Path)
    parser.add_argument("--json", action="store_true", dest="as_json")
    args = parser.parse_args()

    prs = Presentation(args.deck)
    rows = []
    for index, slide in enumerate(prs.slides, 1):
        texts = []
        counts = {
            "pictures": 0,
            "groups": 0,
            "charts": 0,
            "tables": 0,
            "media_shapes": 0,
        }
        image_types = {}
        all_shapes = list(walk_shapes(slide.shapes))
        for shape in all_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                counts["groups"] += 1
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value.replace("\n", " / "))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["pictures"] += 1
                try:
                    ext = shape.image.ext.lower()
                except Exception:
                    ext = "unknown"
                image_types[ext] = image_types.get(ext, 0) + 1
            if getattr(shape, "has_chart", False):
                counts["charts"] += 1
            if getattr(shape, "has_table", False):
                counts["tables"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                counts["media_shapes"] += 1

        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip().replace("\n", " / ")[:160]
        if not title and texts:
            title = texts[0][:160]

        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        rows.append(
            {
                "slide": index,
                "title_or_first_text": title,
                "top_level_shapes": len(slide.shapes),
                "all_shapes_including_groups": len(all_shapes),
                "editable_text_shapes": len(texts),
                **counts,
                "image_types": image_types,
                "has_notes": bool(notes),
                "layout": slide.slide_layout.name,
            }
        )

    result = {
        "file": str(args.deck.resolve()),
        "slides": len(prs.slides),
        "width_inches": round(prs.slide_width / 914400, 3),
        "height_inches": round(prs.slide_height / 914400, 3),
        "inventory": rows,
    }
    if args.as_json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    print(
        f"{result['file']}\n"
        f"slides={result['slides']} size={result['width_inches']}x{result['height_inches']} in"
    )
    for row in rows:
        print(
            f"{row['slide']:02d} | shapes={row['all_shapes_including_groups']:3d} "
            f"text={row['editable_text_shapes']:2d} pic={row['pictures']:2d} "
            f"group={row['groups']:2d} chart={row['charts']} table={row['tables']} "
            f"media={row['media_shapes']} notes={'yes' if row['has_notes'] else 'no'} "
            f"| {row['title_or_first_text']}"
        )


if __name__ == "__main__":
    main()
