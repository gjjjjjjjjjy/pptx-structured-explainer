#!/usr/bin/env python3
"""Create an editable PPTX from a small JSON slide specification."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SHAPES = {
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
}


def color(value: str | None, default: str = "000000") -> RGBColor:
    return RGBColor.from_string((value or default).lstrip("#"))


def box(element: dict) -> tuple:
    return tuple(Inches(float(element[key])) for key in ("x", "y", "w", "h"))


def style_text_frame(frame, element: dict, default_font: str) -> None:
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(float(element.get("margin", 0.08)))
    frame.margin_top = frame.margin_bottom = Inches(float(element.get("margin", 0.08)))
    frame.vertical_anchor = getattr(MSO_ANCHOR, element.get("vertical", "MIDDLE").upper())
    lines = str(element.get("text", "")).split("\n") or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = getattr(PP_ALIGN, element.get("align", "LEFT").upper())
        run = paragraph.add_run()
        run.text = line
        font = run.font
        font.name = element.get("font", default_font)
        font.size = Pt(float(element.get("font_size", 18)))
        font.bold = bool(element.get("bold", False))
        font.italic = bool(element.get("italic", False))
        font.color.rgb = color(element.get("color"), "111827")


def add_element(slide, element: dict, base_dir: Path, default_font: str) -> None:
    kind = element["type"]
    if kind == "text":
        shape = slide.shapes.add_textbox(*box(element))
        style_text_frame(shape.text_frame, element, default_font)
        return
    if kind == "shape":
        shape_type = SHAPES[element.get("shape", "rectangle")]
        shape = slide.shapes.add_shape(shape_type, *box(element))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(element.get("fill"), "FFFFFF")
        shape.line.color.rgb = color(element.get("line"), "D1D5DB")
        shape.line.width = Pt(float(element.get("line_width", 1)))
        if "text" in element:
            style_text_frame(shape.text_frame, element, default_font)
        return
    if kind == "line":
        x, y, w, h = box(element)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y + h)
        line.line.color.rgb = color(element.get("color"), "374151")
        line.line.width = Pt(float(element.get("width", 1.5)))
        return
    if kind == "image":
        image_path = (base_dir / element["path"]).resolve()
        left, top, width, height = box(element)
        with Image.open(image_path) as source:
            image_width, image_height = source.size
        scale = min(width / image_width, height / image_height)
        fitted_width = round(image_width * scale)
        fitted_height = round(image_height * scale)
        fitted_left = left + round((width - fitted_width) / 2)
        fitted_top = top + round((height - fitted_height) / 2)
        slide.shapes.add_picture(str(image_path), fitted_left, fitted_top, fitted_width, fitted_height)
        return
    if kind == "table":
        data = element.get("data", [])
        if not data or not all(isinstance(row, list) for row in data):
            raise ValueError("table data must be a non-empty list of rows")
        columns = max(len(row) for row in data)
        table = slide.shapes.add_table(len(data), columns, *box(element)).table
        for row_index, row in enumerate(data):
            for column_index in range(columns):
                cell = table.cell(row_index, column_index)
                cell.text = str(row[column_index]) if column_index < len(row) else ""
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = element.get("font", default_font)
                        run.font.size = Pt(float(element.get("font_size", 14)))
                        run.font.bold = bool(element.get("header", True) and row_index == 0)
                        run.font.color.rgb = color(element.get("color"), "111827")
        return
    raise ValueError(f"unsupported element type: {kind}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("spec", type=Path)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()

    spec_path = args.spec.expanduser().resolve()
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    prs = Presentation()
    if spec.get("layout", "wide") == "wide":
        prs.slide_width, prs.slide_height = Inches(13.333333), Inches(7.5)
    default_font = spec.get("theme", {}).get("font", "Arial")
    blank = prs.slide_layouts[6]

    for slide_spec in spec.get("slides", []):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = color(slide_spec.get("background"), "FFFFFF")
        for element in slide_spec.get("elements", []):
            add_element(slide, element, spec_path.parent, default_font)

    output = args.output.expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"slides={len(prs.slides)}")
    print(f"output={output}")


if __name__ == "__main__":
    main()
