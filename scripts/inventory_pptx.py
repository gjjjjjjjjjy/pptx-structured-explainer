#!/usr/bin/env python3
"""Print a compact structural inventory for a PPTX/POTX file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from cli_compat import configure_utf8_stdio


configure_utf8_stdio()

def walk_shapes(shapes):
    """Yield top-level shapes and every shape nested inside a group."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck", type=Path)
    parser.add_argument("--json", action="store_true", dest="as_json")
    args = parser.parse_args()

    prs = Presentation(args.deck)
    rows = []
    for index, slide in enumerate(prs.slides, 1):
        texts = []
        pictures = charts = tables = media = 0
        top_level_shapes = len(slide.shapes)
        nested_shapes = 0
        for shape in walk_shapes(slide.shapes):
            nested_shapes += 1
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value.replace("\n", " / "))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                media += 1

        title = texts[0][:120] if texts else ""
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        rows.append(
            {
                "slide": index,
                "title_or_first_text": title,
                "editable_text_shapes": len(texts),
                "pictures": pictures,
                "charts": charts,
                "tables": tables,
                "media_shapes": media,
                "top_level_shapes": top_level_shapes,
                "all_shapes_including_groups": nested_shapes,
                "has_notes": bool(notes),
                "layout": slide.slide_layout.name,
            }
        )

    result = {
        "file": str(args.deck.resolve()),
        "slides": len(prs.slides),
        "width_inches": round(prs.slide_width / 914400, 3),
        "height_inches": round(prs.slide_height / 914400, 3),
        "inventory": rows,
    }
    if args.as_json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    print(
        f"{result['file']}\n"
        f"slides={result['slides']} size={result['width_inches']}x{result['height_inches']} in"
    )
    for row in rows:
        print(
            f"{row['slide']:02d} | text={row['editable_text_shapes']:2d} "
            f"pic={row['pictures']:2d} chart={row['charts']} table={row['tables']} "
            f"media={row['media_shapes']} notes={'yes' if row['has_notes'] else 'no'} "
            f"| {row['title_or_first_text']}"
        )


if __name__ == "__main__":
    main()
