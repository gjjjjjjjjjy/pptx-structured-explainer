#!/usr/bin/env python3
"""Replace exact text in PPTX text frames while preserving unaffected shapes."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def replace_paragraph(paragraph, replacements: dict[str, str]) -> tuple[int, bool]:
    original = "".join(run.text for run in paragraph.runs)
    updated = original
    count = 0
    for source, target in replacements.items():
        occurrences = updated.count(source)
        if occurrences:
            updated = updated.replace(source, target)
            count += occurrences
    if updated == original:
        return 0, False
    if not paragraph.runs:
        paragraph.add_run().text = updated
        return count, False
    paragraph.runs[0].text = updated
    for run in paragraph.runs[1:]:
        run.text = ""
    return count, len(paragraph.runs) > 1


def replace_frame(frame, replacements: dict[str, str]) -> tuple[int, int]:
    count = collapsed = 0
    for paragraph in frame.paragraphs:
        changed, mixed = replace_paragraph(paragraph, replacements)
        count += changed
        collapsed += int(changed > 0 and mixed)
    return count, collapsed


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck", type=Path)
    parser.add_argument("replacements", type=Path, help='JSON object such as {"old": "new"}')
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()

    source = args.deck.expanduser().resolve()
    replacements = json.loads(args.replacements.read_text(encoding="utf-8"))
    if not isinstance(replacements, dict) or not all(isinstance(k, str) and isinstance(v, str) for k, v in replacements.items()):
        parser.error("replacements must be a JSON object containing string keys and values")

    prs = Presentation(source)
    total = collapsed = 0
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                changed, mixed = replace_frame(shape.text_frame, replacements)
                total += changed
                collapsed += mixed
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        changed, mixed = replace_frame(cell.text_frame, replacements)
                        total += changed
                        collapsed += mixed

    output = args.output.expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"replacements={total}")
    print(f"mixed_format_paragraphs_collapsed={collapsed}")
    print(f"output={output}")


if __name__ == "__main__":
    main()
